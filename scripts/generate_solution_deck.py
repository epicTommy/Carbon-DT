from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "artifacts"
OUTPUT_PATH = OUTPUT_DIR / "Carbon_DT_Overview.pptx"
LOGO_PATH = PROJECT_ROOT / "image.png"

BG = RGBColor(247, 249, 252)
NAVY = RGBColor(19, 39, 61)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(14, 116, 144)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(94, 109, 128)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(227, 238, 255)


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullet_block(slide, left, top, width, height, title, bullets, accent):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = LIGHT

    add_textbox(slide, left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), Inches(0.35), title, 18, True, accent)

    box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.55), width - Inches(0.36), height - Inches(0.7))
    frame = box.text_frame
    frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT


def add_header_band(slide, title, subtitle):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_textbox(slide, Inches(0.6), Inches(0.22), Inches(7.8), Inches(0.35), title, 28, True, WHITE)
    add_textbox(slide, Inches(0.62), Inches(0.58), Inches(8.5), Inches(0.2), subtitle, 11, False, RGBColor(210, 223, 241))


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.65), Inches(12.2), Inches(5.9))
    hero.fill.solid()
    hero.fill.fore_color.rgb = WHITE
    hero.line.color.rgb = LIGHT

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(8.75), Inches(0.95), width=Inches(3.1))

    add_textbox(slide, Inches(1.0), Inches(1.2), Inches(6.8), Inches(0.6), "Carbon DT", 28, True, NAVY)
    add_textbox(slide, Inches(1.0), Inches(1.78), Inches(6.5), Inches(0.65), "A building-energy screening MVP that converts raw utility bills into expected usage, diagnostics, operational recommendations, and carbon/compliance outputs.", 17, False, TEXT)

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.55), Inches(2.4), Inches(0.5))
    pill.fill.solid()
    pill.fill.fore_color.rgb = LIGHT
    pill.line.fill.background()
    add_textbox(slide, Inches(1.18), Inches(2.69), Inches(2.1), Inches(0.2), "Uploaded bills + Open-Meteo", 14, True, BLUE)

    add_bullet_block(
        slide,
        Inches(0.95),
        Inches(3.35),
        Inches(3.6),
        Inches(2.5),
        "Problem Being Solved",
        [
            "Utility-bill data is hard to compare month to month without normalization",
            "Owners need a fast first-pass answer before a full engineering audit",
            "Results must be explainable enough for operators, not a black-box model",
        ],
        BLUE,
    )
    add_bullet_block(
        slide,
        Inches(4.85),
        Inches(3.35),
        Inches(3.6),
        Inches(2.5),
        "Core Output",
        [
            "Monthly electricity and gas records normalized to kWh and therms",
            "Expected-vs-actual energy baseline adjusted for weather and season",
            "Diagnostics, deterministic recommendations, emissions, and LL97 screening",
        ],
        TEAL,
    )
    add_bullet_block(
        slide,
        Inches(8.75),
        Inches(3.35),
        Inches(3.2),
        Inches(2.5),
        "Current Delivery Mode",
        [
            "Public Streamlit web app with CSV upload",
            "Historical weather fetched from Open-Meteo by ZIP or address",
            "No embedded customer files, secrets, or LLM dependency",
        ],
        NAVY,
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_header_band(slide, "How The Solution Works", "Concrete data flow from uploaded bills to expected usage and recommended actions")

    boxes = [
        (
            "1. Utility Bill Ingestion",
            [
                "CSV input requires: billing_start, billing_end, utility_type, usage, usage_unit, cost",
                "Electricity is normalized to kWh and gas to therms",
                "Billing periods spanning multiple months are prorated into monthly records",
            ],
            BLUE,
        ),
        (
            "2. Weather And Monthly Features",
            [
                "Open-Meteo returns historical daily mean temperature for the billing window",
                "The app aggregates to monthly avg_temp and calculates HDD/CDD at a 65F base",
                "It also creates month index, season flags, and energy-per-sqft features",
            ],
            TEAL,
        ),
        (
            "3. Expected Energy Calculation",
            [
                "Electricity and gas are modeled separately because their drivers differ",
                "If history is sufficient, a simple sklearn linear regression is fit for explainability",
                "If history is sparse, the app falls back to a weather-normalized heuristic baseline",
            ],
            NAVY,
        ),
        (
            "4. Diagnostics To Action",
            [
                "Rules detect patterns such as high electricity, winter gas spike, baseload, and negative trend",
                "Recommendations come from a fixed library mapped to each diagnostic",
                "Ranking uses a transparent weighted formula based on severity, confidence, category, and difficulty",
            ],
            BLUE,
        ),
    ]

    positions = [
        (Inches(0.8), Inches(1.45)),
        (Inches(6.95), Inches(1.45)),
        (Inches(0.8), Inches(4.0)),
        (Inches(6.95), Inches(4.0)),
    ]
    for (title, bullets, accent), (left, top) in zip(boxes, positions):
        add_bullet_block(slide, left, top, Inches(5.55), Inches(2.15), title, bullets, accent)

    arrow1 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(5.95), Inches(2.32), Inches(0.6), Inches(0.45))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = LIGHT
    arrow1.line.fill.background()
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(5.95), Inches(4.9), Inches(0.6), Inches(0.45))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = LIGHT
    arrow2.line.fill.background()

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(6.25), Inches(11.5), Inches(0.75))
    note.fill.solid()
    note.fill.fore_color.rgb = WHITE
    note.line.color.rgb = LIGHT
    add_textbox(
        slide,
        Inches(1.15),
        Inches(6.43),
        Inches(11.0),
        Inches(0.35),
        "Example baseline logic: predicted electricity intensity uses month index, CDD, HDD, summer/winter flags; gas uses HDD, CDD, heating season, and winter flags. Output is predicted intensity x floor area.",
        13,
        False,
        MUTED,
        PP_ALIGN.CENTER,
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_header_band(slide, "Current MVP Status And Next Steps", "Where the product is credible today and what should improve before larger rollout")

    add_bullet_block(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(3.85),
        Inches(4.6),
        "What Is Already Working",
        [
            "Upload-driven screening workflow with validation messages",
            "Weather-aware expected usage for electricity and gas",
            "Deterministic diagnostics and recommendation ranking",
            "Annual emissions calculation and NYC LL97 penalty estimate",
            "Downloadable PDF report for audit-style output",
        ],
        BLUE,
    )
    add_bullet_block(
        slide,
        Inches(4.75),
        Inches(1.55),
        Inches(3.85),
        Inches(4.6),
        "Important Assumptions",
        [
            "This is a single-building MVP, not a portfolio platform yet",
            "Weather quality depends on correct ZIP/address geocoding",
            "Savings and carbon ranges are heuristic placeholders",
            "Expected usage is for screening and triage, not investment-grade M&V",
        ],
        TEAL,
    )
    add_bullet_block(
        slide,
        Inches(8.7),
        Inches(1.55),
        Inches(3.85),
        Inches(4.6),
        "Recommended Next Moves",
        [
            "Deploy and validate with real customer bills",
            "Measure cold-start and weather-fetch latency with production telemetry",
            "Calibrate diagnostics and recommendation heuristics with audit outcomes",
            "Then add API separation, stronger UI, and portfolio support",
        ],
        NAVY,
    )

    footer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(6.45), Inches(11.9), Inches(0.55))
    footer.fill.solid()
    footer.fill.fore_color.rgb = NAVY
    footer.line.fill.background()
    add_textbox(
        slide,
        Inches(1.25),
        Inches(6.62),
        Inches(11.3),
        Inches(0.2),
        "Positioning: strong for fast building-level audit screening and operator guidance; not yet a substitute for full engineering analysis or enterprise workflow software.",
        15,
        False,
        WHITE,
        PP_ALIGN.CENTER,
    )

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    build_deck()
